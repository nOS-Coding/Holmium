"""Document reading tools — PDF, DOCX, XLSX, PPTX extraction to markdown."""

import os
from pathlib import Path
from typing import Any, Dict, List, Optional

import httpx

from tools.registry import register_tool

VLLM_SOCKET = "/run/holmium/vllm.sock"


def _call_vllm(prompt: str) -> str:
    transport = httpx.HTTPTransport(uds=VLLM_SOCKET)
    payload = {
        "model": "Qwen3.6-35B-A3B-AWQ",
        "messages": [{"role": "user", "content": prompt}],
        "max_tokens": 2048,
        "temperature": 0.2,
    }
    with httpx.Client(transport=transport, timeout=120) as client:
        resp = client.post("http://localhost/v1/chat/completions", json=payload)
        resp.raise_for_status()
        return resp.json()["choices"][0]["message"]["content"]


def _read_pdf(path: str) -> str:
    import fitz
    doc = fitz.open(path)
    lines: List[str] = []
    for page in doc:
        lines.append(page.get_text())
    doc.close()
    return "\n".join(lines)


def _read_docx(path: str) -> str:
    import docx
    doc = docx.Document(path)
    lines = [p.text for p in doc.paragraphs]
    return "\n".join(lines)


def _read_xlsx(path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    lines: List[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"## Sheet: {sheet_name}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            lines.append(" | ".join(cells))
    wb.close()
    return "\n".join(lines)


def _read_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    lines: List[str] = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
        lines.append("")
    return "\n".join(lines)


@register_tool(
    "doc_read",
    "Read a document file (PDF/DOCX/XLSX/PPTX) and extract content as markdown.",
    params_schema={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "Absolute path to the document file",
            },
        },
        "required": ["path"],
    },
)
def doc_read(path: str) -> str:
    """Auto-detect document type by extension and extract to markdown."""
    if not os.path.isfile(path):
        return f"Error: file not found at {path}"

    try:
        ext = Path(path).suffix.lower()
        if ext == ".pdf":
            return _read_pdf(path)
        elif ext == ".docx":
            return _read_docx(path)
        elif ext == ".xlsx":
            return _read_xlsx(path)
        elif ext == ".pptx":
            return _read_pptx(path)
        else:
            return f"Unsupported file type: {ext}. Supported: .pdf, .docx, .xlsx, .pptx"
    except Exception as e:
        return f"Error reading document: {e}"


@register_tool(
    "doc_summarize",
    "Read a document and return an AI-generated summary via vLLM.",
    params_schema={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "Absolute path to the document file",
            },
        },
        "required": ["path"],
    },
)
def doc_summarize(path: str) -> str:
    """Read a document and pass the content to vLLM for summarization."""
    content = doc_read(path)
    if content.startswith("Error"):
        return content

    prompt = (
        "Summarize the following document content. "
        "Provide: key topics, main points, and any important data or conclusions. "
        "Keep the summary concise.\n\n"
        f"---DOCUMENT CONTENT---\n{content[:15000]}"
    )
    try:
        return _call_vllm(prompt)
    except Exception as e:
        return f"Error summarizing document: {e}"
